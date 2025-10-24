"""
Advanced Document Processing Service
Supports multiple formats with intelligent content extraction
"""

import os
import asyncio
from typing import Dict, Any, Optional, List
import logging
from pathlib import Path
import mimetypes

# Document processing libraries
import csv
import json
import xml.etree.ElementTree as ET
import zipfile
import markdown
import yaml

# Optional document libraries
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import eml_parser
    EML_AVAILABLE = True
except ImportError:
    EML_AVAILABLE = False

try:
    import rarfile
    RAR_AVAILABLE = True
except ImportError:
    RAR_AVAILABLE = False

# OCR and image processing
try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# Audio/Video metadata
try:
    import mutagen
    AUDIO_METADATA_AVAILABLE = True
except ImportError:
    AUDIO_METADATA_AVAILABLE = False

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self):
        self.supported_formats = {
            # Text formats
            '.txt': self._process_text,
            '.md': self._process_markdown,
            '.rst': self._process_text,
            
            # PDF
            '.pdf': self._process_pdf,
            
            # Microsoft Office
            '.docx': self._process_docx,
            '.doc': self._process_doc,
            '.xlsx': self._process_xlsx,
            '.xls': self._process_excel,
            '.pptx': self._process_pptx,
            '.ppt': self._process_powerpoint,
            
            # Data formats
            '.csv': self._process_csv,
            '.json': self._process_json,
            '.xml': self._process_xml,
            '.yaml': self._process_yaml,
            '.yml': self._process_yaml,
            
            # Web formats
            '.html': self._process_html,
            '.htm': self._process_html,
            
            # Email
            '.eml': self._process_email,
            '.msg': self._process_email,
            
            # Archives
            '.zip': self._process_zip,
            '.rar': self._process_rar,
            
            # Images (with OCR)
            '.png': self._process_image,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.gif': self._process_image,
            '.bmp': self._process_image,
            '.tiff': self._process_image,
            
            # Code files
            '.py': self._process_code,
            '.js': self._process_code,
            '.ts': self._process_code,
            '.java': self._process_code,
            '.cpp': self._process_code,
            '.c': self._process_code,
            '.go': self._process_code,
            '.rs': self._process_code,
            '.php': self._process_code,
            '.rb': self._process_code,
            '.sql': self._process_code,
        }
    
    async def process_document(self, file_path: str) -> Dict[str, Any]:
        """Process document and extract content"""
        try:
            file_path = Path(file_path)
            file_extension = file_path.suffix.lower()
            
            if file_extension not in self.supported_formats:
                return {
                    "success": False,
                    "error": f"Unsupported file format: {file_extension}",
                    "supported_formats": list(self.supported_formats.keys())
                }
            
            # Get file metadata
            metadata = await self._get_file_metadata(file_path)
            
            # Process content based on file type
            processor = self.supported_formats[file_extension]
            content_data = await processor(file_path)
            
            return {
                "success": True,
                "file_path": str(file_path),
                "file_type": file_extension,
                "metadata": metadata,
                **content_data
            }
            
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {e}")
            return {
                "success": False,
                "error": str(e),
                "file_path": str(file_path)
            }
    
    async def _get_file_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract file metadata"""
        stat = file_path.stat()
        mime_type, _ = mimetypes.guess_type(str(file_path))
        
        return {
            "filename": file_path.name,
            "size": stat.st_size,
            "created": stat.st_ctime,
            "modified": stat.st_mtime,
            "mime_type": mime_type,
            "extension": file_path.suffix.lower()
        }
    
    # Text processing methods
    async def _process_text(self, file_path: Path) -> Dict[str, Any]:
        """Process plain text files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        return {
            "content": content,
            "word_count": len(content.split()),
            "char_count": len(content),
            "line_count": len(content.splitlines())
        }
    
    async def _process_markdown(self, file_path: Path) -> Dict[str, Any]:
        """Process Markdown files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            md_content = f.read()
        
        # Convert to HTML for better parsing
        html_content = markdown.markdown(md_content, extensions=['meta', 'toc'])
        
        return {
            "content": md_content,
            "html_content": html_content,
            "word_count": len(md_content.split()),
            "char_count": len(md_content)
        }
    
    async def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Process PDF files"""
        content = ""
        page_count = 0
        
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            page_count = len(pdf_reader.pages)
            
            for page in pdf_reader.pages:
                content += page.extract_text() + "\n"
        
        return {
            "content": content,
            "page_count": page_count,
            "word_count": len(content.split()),
            "char_count": len(content)
        }
    
    async def _process_docx(self, file_path: Path) -> Dict[str, Any]:
        """Process DOCX files"""
        doc = docx.Document(file_path)
        content = ""
        
        for paragraph in doc.paragraphs:
            content += paragraph.text + "\n"
        
        # Extract tables
        tables_content = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables_content.append(table_data)
        
        return {
            "content": content,
            "tables": tables_content,
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(doc.tables),
            "word_count": len(content.split())
        }
    
    async def _process_xlsx(self, file_path: Path) -> Dict[str, Any]:
        """Process Excel files"""
        workbook = openpyxl.load_workbook(file_path)
        sheets_data = {}
        content = ""
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_data = []
            
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    sheet_data.append([str(cell) if cell is not None else "" for cell in row])
                    content += " ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
            
            sheets_data[sheet_name] = sheet_data
        
        return {
            "content": content,
            "sheets": sheets_data,
            "sheet_count": len(workbook.sheetnames),
            "word_count": len(content.split())
        }
    
    async def _process_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Process PowerPoint files"""
        prs = Presentation(file_path)
        content = ""
        slides_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            
            slides_content.append({
                "slide_number": i + 1,
                "content": slide_text
            })
            content += slide_text + "\n"
        
        return {
            "content": content,
            "slides": slides_content,
            "slide_count": len(prs.slides),
            "word_count": len(content.split())
        }
    
    async def _process_csv(self, file_path: Path) -> Dict[str, Any]:
        """Process CSV files"""
        content = ""
        rows_data = []
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            csv_reader = csv.reader(f)
            headers = next(csv_reader, None)
            
            for row in csv_reader:
                rows_data.append(row)
                content += " ".join(row) + "\n"
        
        return {
            "content": content,
            "headers": headers,
            "rows": rows_data,
            "row_count": len(rows_data),
            "column_count": len(headers) if headers else 0
        }
    
    async def _process_json(self, file_path: Path) -> Dict[str, Any]:
        """Process JSON files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Convert to readable text
        content = json.dumps(data, indent=2)
        
        return {
            "content": content,
            "json_data": data,
            "structure": self._analyze_json_structure(data),
            "word_count": len(content.split())
        }
    
    async def _process_xml(self, file_path: Path) -> Dict[str, Any]:
        """Process XML files"""
        tree = ET.parse(file_path)
        root = tree.getroot()
        
        content = ET.tostring(root, encoding='unicode')
        
        return {
            "content": content,
            "root_tag": root.tag,
            "element_count": len(list(root.iter())),
            "word_count": len(content.split())
        }
    
    async def _process_yaml(self, file_path: Path) -> Dict[str, Any]:
        """Process YAML files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = yaml.safe_load(f)
        
        content = yaml.dump(data, default_flow_style=False)
        
        return {
            "content": content,
            "yaml_data": data,
            "word_count": len(content.split())
        }
    
    async def _process_image(self, file_path: Path) -> Dict[str, Any]:
        """Process images with OCR"""
        if not OCR_AVAILABLE:
            return {
                "content": "",
                "error": "OCR not available. Install pytesseract and PIL."
            }
        
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            
            return {
                "content": text,
                "image_size": image.size,
                "image_mode": image.mode,
                "word_count": len(text.split()),
                "ocr_confidence": "Available with pytesseract"
            }
        except Exception as e:
            return {
                "content": "",
                "error": f"OCR failed: {str(e)}"
            }
    
    async def _process_code(self, file_path: Path) -> Dict[str, Any]:
        """Process code files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Basic code analysis
        lines = content.splitlines()
        non_empty_lines = [line for line in lines if line.strip()]
        comment_lines = [line for line in lines if line.strip().startswith(('#', '//', '/*', '*', '--'))]
        
        return {
            "content": content,
            "total_lines": len(lines),
            "code_lines": len(non_empty_lines),
            "comment_lines": len(comment_lines),
            "language": file_path.suffix[1:],  # Remove the dot
            "word_count": len(content.split())
        }
    
    # Placeholder methods for formats requiring additional libraries
    async def _process_doc(self, file_path: Path) -> Dict[str, Any]:
        """Process legacy DOC files"""
        return {"content": "", "error": "Legacy DOC format requires python-docx2txt"}
    
    async def _process_excel(self, file_path: Path) -> Dict[str, Any]:
        """Process legacy Excel files"""
        return {"content": "", "error": "Legacy XLS format requires xlrd"}
    
    async def _process_powerpoint(self, file_path: Path) -> Dict[str, Any]:
        """Process legacy PowerPoint files"""
        return {"content": "", "error": "Legacy PPT format requires additional libraries"}
    
    async def _process_html(self, file_path: Path) -> Dict[str, Any]:
        """Process HTML files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Basic HTML parsing (could be enhanced with BeautifulSoup)
        return {
            "content": content,
            "html_content": content,
            "word_count": len(content.split())
        }
    
    async def _process_email(self, file_path: Path) -> Dict[str, Any]:
        """Process email files"""
        return {"content": "", "error": "Email processing requires eml-parser"}
    
    async def _process_zip(self, file_path: Path) -> Dict[str, Any]:
        """Process ZIP archives"""
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                file_list = zip_file.namelist()
                return {
                    "content": f"ZIP archive containing {len(file_list)} files",
                    "files": file_list,
                    "file_count": len(file_list)
                }
        except Exception as e:
            return {"content": "", "error": f"Failed to process ZIP: {str(e)}"}
    
    async def _process_rar(self, file_path: Path) -> Dict[str, Any]:
        """Process RAR archives"""
        return {"content": "", "error": "RAR processing requires rarfile library"}
    
    def _analyze_json_structure(self, data: Any, depth: int = 0) -> Dict[str, Any]:
        """Analyze JSON structure"""
        if isinstance(data, dict):
            return {
                "type": "object",
                "keys": list(data.keys()),
                "depth": depth,
                "size": len(data)
            }
        elif isinstance(data, list):
            return {
                "type": "array",
                "length": len(data),
                "depth": depth,
                "item_types": list(set(type(item).__name__ for item in data[:10]))  # Sample first 10
            }
        else:
            return {
                "type": type(data).__name__,
                "depth": depth
            }

# Global service instance
document_processor = DocumentProcessor()
